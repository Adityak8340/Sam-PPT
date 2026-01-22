"""
PowerPoint Generator Module
Creates professional presentations with images and web content
"""

import os
from typing import Optional, Any, cast
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor  # type: ignore
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # type: ignore
from pptx.enum.shapes import MSO_SHAPE  # type: ignore
from PIL import Image

from web_scraper import WebPageData


class PPTGenerator:
    """Generates PowerPoint presentations from web search results"""
    
    # Professional color scheme
    COLORS = {
        'primary': RGBColor(0x1a, 0x73, 0xe8),      # Blue
        'secondary': RGBColor(0x34, 0xa8, 0x53),    # Green
        'accent': RGBColor(0xfb, 0xbc, 0x04),       # Yellow
        'dark': RGBColor(0x20, 0x2a, 0x44),         # Dark blue
        'light': RGBColor(0xf8, 0xf9, 0xfa),        # Light gray
        'text': RGBColor(0x33, 0x33, 0x33),         # Dark gray
        'white': RGBColor(0xff, 0xff, 0xff),        # White
    }
    
    def __init__(self) -> None:
        self.prs: Any = Presentation()
        # Set slide dimensions (16:9 widescreen)
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
    
    def create_presentation(self, input_image_path: str, 
                           search_results: list[dict],
                           web_data: list[WebPageData],
                           output_path: str) -> str:
        """
        Create a complete presentation
        
        Args:
            input_image_path: Path to the original input image
            search_results: List of search results with title and url
            web_data: List of WebPageData objects
            output_path: Path to save the presentation
            
        Returns:
            Path to the saved presentation
        """
        # Create title slide
        self._add_title_slide(input_image_path)
        
        # Create search results summary slide
        self._add_summary_slide(search_results)
        
        # Create detailed slides for each website
        for i, (result, data) in enumerate(zip(search_results, web_data)):
            self._add_website_slide(i + 1, result, data)
        
        # Create conclusion slide
        self._add_conclusion_slide(len(search_results))
        
        # Save presentation
        self.prs.save(output_path)
        return output_path
    
    def _add_title_slide(self, image_path: str):
        """Add title slide with the input image"""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add background shape
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.COLORS['dark']
        background.line.fill.background()
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(12.333), Inches(1)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "Image Search Results"
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['white']
        title_para.alignment = PP_ALIGN.CENTER
        
        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(12.333), Inches(0.5)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = "Reverse Image Search Analysis"
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = self.COLORS['light']
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Add the input image
        try:
            # Resize image if needed
            img = Image.open(image_path)
            max_width, max_height = 500, 350
            img.thumbnail((max_width, max_height), Image.Resampling.LANCZOS)
            
            # Save to bytes
            img_bytes = BytesIO()
            img_format = 'PNG' if image_path.lower().endswith('.png') else 'JPEG'
            img.save(img_bytes, format=img_format)
            img_bytes.seek(0)
            
            # Calculate centered position
            img_width = Inches(img.width / 96)  # Convert pixels to inches (96 DPI)
            img_height = Inches(img.height / 96)
            left = (self.prs.slide_width - img_width) / 2
            top = Inches(2.5)
            
            # Add image with border
            border_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left - Inches(0.15), top - Inches(0.15),
                img_width + Inches(0.3), img_height + Inches(0.3)
            )
            border_shape.fill.solid()
            border_shape.fill.fore_color.rgb = self.COLORS['white']
            border_shape.line.fill.background()
            
            slide.shapes.add_picture(img_bytes, left, top, img_width, img_height)
            
        except Exception:
            # Add placeholder text if image fails
            error_box = slide.shapes.add_textbox(
                Inches(4), Inches(3),
                Inches(5), Inches(1)
            )
            error_frame = error_box.text_frame
            error_para = error_frame.paragraphs[0]
            error_para.text = f"[Input Image]\n{os.path.basename(image_path)}"
            error_para.font.size = Pt(18)
            error_para.font.color.rgb = self.COLORS['light']
            error_para.alignment = PP_ALIGN.CENTER
        
        # Add label
        label_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.8),
            Inches(12.333), Inches(0.5)
        )
        label_frame = label_box.text_frame
        label_para = label_frame.paragraphs[0]
        label_para.text = "Input Image for Analysis"
        label_para.font.size = Pt(16)
        label_para.font.color.rgb = self.COLORS['accent']
        label_para.alignment = PP_ALIGN.CENTER
    
    def _add_summary_slide(self, search_results: list[dict]):
        """Add summary slide with all found websites"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add header background
        header_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, Inches(1.5)
        )
        header_bg.fill.solid()
        header_bg.fill.fore_color.rgb = self.COLORS['primary']
        header_bg.line.fill.background()
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4),
            Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = f"Top {len(search_results)} Matching Websites"
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['white']
        
        # Add website list
        y_position = Inches(1.8)
        for i, result in enumerate(search_results):
            # Number badge
            badge = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.5), y_position,
                Inches(0.5), Inches(0.5)
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = self.COLORS['secondary']
            badge.line.fill.background()
            
            # Number text
            badge.text_frame.paragraphs[0].text = str(i + 1)
            badge.text_frame.paragraphs[0].font.size = Pt(16)
            badge.text_frame.paragraphs[0].font.bold = True
            badge.text_frame.paragraphs[0].font.color.rgb = self.COLORS['white']
            badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            badge.text_frame.word_wrap = False
            
            # Website title
            title_box = slide.shapes.add_textbox(
                Inches(1.2), y_position,
                Inches(11), Inches(0.4)
            )
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = result.get('title', 'Unknown')[:80]
            title_para.font.size = Pt(18)
            title_para.font.bold = True
            title_para.font.color.rgb = self.COLORS['dark']
            
            # URL
            url_box = slide.shapes.add_textbox(
                Inches(1.2), y_position + Inches(0.35),
                Inches(11), Inches(0.3)
            )
            url_frame = url_box.text_frame
            url_para = url_frame.paragraphs[0]
            url_para.text = result.get('url', '')[:100]
            url_para.font.size = Pt(12)
            url_para.font.color.rgb = self.COLORS['primary']
            
            y_position += Inches(0.9)
    
    def _add_website_slide(self, index: int, _result: dict, data: WebPageData):
        """Add detailed slide for a single website"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Header with site number
        header_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, Inches(1.2)
        )
        header_bg.fill.solid()
        header_bg.fill.fore_color.rgb = self.COLORS['primary']
        header_bg.line.fill.background()
        
        # Site number and title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = f"#{index}: {data.title[:60]}"
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['white']
        
        # Two-column layout
        # Left column: Screenshot
        if data.screenshot_path and os.path.exists(data.screenshot_path):
            try:
                img = Image.open(data.screenshot_path)
                # Scale to fit
                max_width, max_height = 500, 350
                img.thumbnail((max_width, max_height), Image.Resampling.LANCZOS)
                
                img_bytes = BytesIO()
                img.save(img_bytes, format='PNG')
                img_bytes.seek(0)
                
                img_width = Inches(img.width / 96)
                img_height = Inches(img.height / 96)
                
                # Add border
                border = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.4), Inches(1.5),
                    img_width + Inches(0.2), img_height + Inches(0.2)
                )
                border.fill.solid()
                border.fill.fore_color.rgb = self.COLORS['light']
                border.line.color.rgb = RGBColor(0xdd, 0xdd, 0xdd)
                
                slide.shapes.add_picture(
                    img_bytes, Inches(0.5), Inches(1.6),
                    img_width, img_height
                )
                
            except Exception:
                self._add_placeholder(slide, Inches(0.5), Inches(1.6), 
                                     Inches(5), Inches(3.5), "Screenshot unavailable")
        else:
            self._add_placeholder(slide, Inches(0.5), Inches(1.6), 
                                 Inches(5), Inches(3.5), "Screenshot unavailable")
        
        # Right column: Content
        content_x = Inches(6.5)
        content_y = Inches(1.4)
        
        # URL
        url_box = slide.shapes.add_textbox(
            content_x, content_y,
            Inches(6.3), Inches(0.4)
        )
        url_frame = url_box.text_frame
        url_para = url_frame.paragraphs[0]
        url_para.text = f"🔗 {data.url[:70]}"
        url_para.font.size = Pt(11)
        url_para.font.color.rgb = self.COLORS['primary']
        
        content_y += Inches(0.5)
        
        # Meta description
        if data.meta_description:
            desc_label = slide.shapes.add_textbox(
                content_x, content_y,
                Inches(6.3), Inches(0.3)
            )
            desc_label.text_frame.paragraphs[0].text = "Description:"
            desc_label.text_frame.paragraphs[0].font.size = Pt(12)
            desc_label.text_frame.paragraphs[0].font.bold = True
            desc_label.text_frame.paragraphs[0].font.color.rgb = self.COLORS['dark']
            
            content_y += Inches(0.3)
            
            desc_box = slide.shapes.add_textbox(
                content_x, content_y,
                Inches(6.3), Inches(0.8)
            )
            desc_frame = desc_box.text_frame
            desc_frame.word_wrap = True
            desc_para = desc_frame.paragraphs[0]
            desc_para.text = data.meta_description[:250]
            desc_para.font.size = Pt(10)
            desc_para.font.color.rgb = self.COLORS['text']
            
            content_y += Inches(0.9)
        
        # Key headings
        if data.headings:
            heading_label = slide.shapes.add_textbox(
                content_x, content_y,
                Inches(6.3), Inches(0.3)
            )
            heading_label.text_frame.paragraphs[0].text = "Key Topics:"
            heading_label.text_frame.paragraphs[0].font.size = Pt(12)
            heading_label.text_frame.paragraphs[0].font.bold = True
            heading_label.text_frame.paragraphs[0].font.color.rgb = self.COLORS['dark']
            
            content_y += Inches(0.35)
            
            headings_text = "\n".join([f"• {h[:60]}" for h in data.headings[:5]])
            headings_box = slide.shapes.add_textbox(
                content_x, content_y,
                Inches(6.3), Inches(1.5)
            )
            headings_frame = headings_box.text_frame
            headings_frame.word_wrap = True
            headings_para = headings_frame.paragraphs[0]
            headings_para.text = headings_text
            headings_para.font.size = Pt(10)
            headings_para.font.color.rgb = self.COLORS['text']
            
            content_y += Inches(1.6)
        
        # Content preview
        if data.paragraphs:
            content_label = slide.shapes.add_textbox(
                content_x, content_y,
                Inches(6.3), Inches(0.3)
            )
            content_label.text_frame.paragraphs[0].text = "Content Preview:"
            content_label.text_frame.paragraphs[0].font.size = Pt(12)
            content_label.text_frame.paragraphs[0].font.bold = True
            content_label.text_frame.paragraphs[0].font.color.rgb = self.COLORS['dark']
            
            content_y += Inches(0.35)
            
            preview_text = data.paragraphs[0][:300] + "..." if len(data.paragraphs[0]) > 300 else data.paragraphs[0]
            preview_box = slide.shapes.add_textbox(
                content_x, content_y,
                Inches(6.3), Inches(1.5)
            )
            preview_frame = preview_box.text_frame
            preview_frame.word_wrap = True
            preview_para = preview_frame.paragraphs[0]
            preview_para.text = preview_text
            preview_para.font.size = Pt(9)
            preview_para.font.color.rgb = self.COLORS['text']
        
        # Error indicator if any
        if data.error:
            error_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(6.8),
                Inches(12), Inches(0.4)
            )
            error_frame = error_box.text_frame
            error_para = error_frame.paragraphs[0]
            error_para.text = f"⚠️ Note: {data.error}"
            error_para.font.size = Pt(10)
            error_para.font.color.rgb = RGBColor(0xdc, 0x35, 0x45)
    
    def _add_placeholder(self, slide, left, top, width, height, text: str):
        """Add a placeholder shape with text"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.COLORS['light']
        shape.line.color.rgb = RGBColor(0xcc, 0xcc, 0xcc)
        
        shape.text_frame.paragraphs[0].text = text
        shape.text_frame.paragraphs[0].font.size = Pt(14)
        shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        shape.text_frame.word_wrap = True
    
    def _add_conclusion_slide(self, total_sites: int):
        """Add conclusion slide"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Background
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.COLORS['dark']
        background.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5),
            Inches(12.333), Inches(1)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "Analysis Complete"
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['white']
        title_para.alignment = PP_ALIGN.CENTER
        
        # Summary
        summary_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.8),
            Inches(12.333), Inches(1)
        )
        summary_frame = summary_box.text_frame
        summary_para = summary_frame.paragraphs[0]
        summary_para.text = f"Found and analyzed {total_sites} matching websites"
        summary_para.font.size = Pt(24)
        summary_para.font.color.rgb = self.COLORS['light']
        summary_para.alignment = PP_ALIGN.CENTER
        
        # Footer
        footer_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.5),
            Inches(12.333), Inches(0.5)
        )
        footer_frame = footer_box.text_frame
        footer_para = footer_frame.paragraphs[0]
        footer_para.text = "Generated by Sam-PPT • Image Search Analysis Tool"
        footer_para.font.size = Pt(14)
        footer_para.font.color.rgb = self.COLORS['accent']
        footer_para.alignment = PP_ALIGN.CENTER


def create_presentation(input_image_path: str,
                       search_results: list[dict],
                       web_data: list,
                       output_path: str) -> str:
    """
    Main function to create a presentation
    
    Args:
        input_image_path: Path to the original image
        search_results: List of search results
        web_data: List of WebPageData objects
        output_path: Output path for the presentation
        
    Returns:
        Path to the saved presentation
    """
    generator = PPTGenerator()
    return generator.create_presentation(
        input_image_path, search_results, web_data, output_path
    )
